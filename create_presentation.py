"""
DeusExMachina - World-Class Presentation Generator
Style: Dark Tech Horror (Custom Brand)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ─── DESIGN SYSTEM ────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x09, 0x09, 0x0F)   # Near black
BG_CARD       = RGBColor(0x14, 0x14, 0x1E)   # Card surface
BG_SECTION    = RGBColor(0xC7, 0x3E, 0x1D)   # Blood red (section divider)
BG_NIGHT      = RGBColor(0x0A, 0x14, 0x28)   # Deep night blue

C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_PARCHMENT   = RGBColor(0xE8, 0xD5, 0xB7)   # Warm off-white
C_RED         = RGBColor(0xC7, 0x3E, 0x1D)   # Blood red
C_GOLD        = RGBColor(0xFF, 0xD7, 0x00)   # Gold accent
C_TEAL        = RGBColor(0x2E, 0xCC, 0xCC)   # Teal/cyan
C_SLATE       = RGBColor(0x7A, 0x8F, 0xA6)   # Muted slate
C_DIMWHITE    = RGBColor(0xB0, 0xBA, 0xC7)   # Dim white for body
C_GREEN       = RGBColor(0x3D, 0xBF, 0x7A)   # Success green
C_ORANGE      = RGBColor(0xF0, 0x8C, 0x00)   # Warning orange
C_PURPLE      = RGBColor(0x9B, 0x59, 0xB6)   # Purple

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── HELPERS ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # totally blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=24, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, font_name="Malgun Gothic"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_multiline_text(slide, lines, x, y, w, h,
                       font_size=20, bold=False, color=C_WHITE,
                       align=PP_ALIGN.LEFT, line_spacing=1.15, font_name="Malgun Gothic"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    for i, (line_text, line_size, line_bold, line_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_size)
        run.font.bold = line_bold
        run.font.color.rgb = line_color
        run.font.name = font_name
    return txBox

def add_divider(slide, x, y, w, color, thickness=0.02):
    shape = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ─── SLIDE BUILDERS ───────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1: Hero Title"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    # Left blood-red accent bar
    add_rect(s, 0, 0, 0.5, 7.5, C_RED)

    # Subtle grid texture effect - horizontal lines
    for i in range(1, 8):
        add_rect(s, 0.5, i * 0.95, 12.83, 0.005, RGBColor(0x20, 0x20, 0x30))

    # Top label
    add_text(s, "LIKELION × DEUS EX MACHINA", 1.0, 0.4, 11, 0.5,
             font_size=13, color=C_RED, bold=True, font_name="Malgun Gothic")

    # Main title - very large
    add_text(s, "DeusExMachina", 1.0, 1.2, 11, 2.0,
             font_size=80, bold=True, color=C_WHITE, font_name="Malgun Gothic")

    # Subtitle
    add_text(s, "AI 기반 인터랙티브 내러티브 게임 엔진", 1.0, 3.1, 11, 0.8,
             font_size=32, bold=False, color=C_PARCHMENT, font_name="Malgun Gothic")

    # Description
    add_text(s,
        "LLM이 살아 숨쉬는 NPC를 만든다 — Generative Agents × 심리 공포 × 탈출 어드벤처",
        1.0, 3.85, 11, 0.6,
        font_size=18, color=C_SLATE, italic=True, font_name="Malgun Gothic")

    add_divider(s, 1.0, 4.7, 6, C_RED, 0.04)

    # Bottom info
    add_text(s, "2025  ·  Python / FastAPI / OpenAI GPT-4  ·  Coraline Scenario v3", 1.0, 5.0, 11, 0.4,
             font_size=14, color=C_SLATE, font_name="Malgun Gothic")

    # Decorative corner
    add_rect(s, 12.33, 6.5, 1.0, 0.04, C_GOLD)
    add_text(s, "인형의 집에서 탈출하라", 9.5, 6.6, 3.5, 0.5,
             font_size=14, color=C_GOLD, italic=True,
             align=PP_ALIGN.RIGHT, font_name="Malgun Gothic")
    return s


def slide_section_divider(prs, number, title, subtitle):
    """Section divider slide - blood red background"""
    s = blank_slide(prs)
    fill_bg(s, C_RED)

    # Dark overlay left portion
    add_rect(s, 0, 0, 4, 7.5, RGBColor(0x8B, 0x1A, 0x0A))

    # Section number
    add_text(s, f"{number:02d}", 0.5, 1.8, 3, 2.5,
             font_size=140, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER, font_name="Malgun Gothic")

    # Vertical divider
    add_rect(s, 4.1, 0.5, 0.04, 6.5, C_WHITE)

    # Title
    add_text(s, title, 4.6, 2.5, 8.3, 1.2,
             font_size=52, bold=True, color=C_WHITE, font_name="Malgun Gothic")

    # Subtitle
    add_text(s, subtitle, 4.6, 3.8, 8.3, 1.0,
             font_size=24, color=RGBColor(0xFF, 0xDD, 0xCC), italic=True, font_name="Malgun Gothic")
    return s


def slide_project_background(prs):
    """Slide 3: Why we built this"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    # Header bar
    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "프로젝트 배경", 0.6, 0.25, 8, 0.7,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "왜 이걸 만들었나?", 0.6, 0.72, 5, 0.4,
             font_size=16, color=C_RED, font_name="Malgun Gothic")

    # Three problem cards
    cards = [
        ("😴", "기존 텍스트 RPG의 한계",
         "정해진 분기, 고정된 대사\n플레이어는 '읽는' 것에 그침\n몰입감 부재"),
        ("🤖", "LLM의 폭발적 성장",
         "GPT-4급 모델의 등장\n자연어 이해·생성 능력 혁신\n실시간 대화 NPC 가능성"),
        ("🎭", "Generative Agents 논문",
         "Park et al. 2023\n기억·계획·반성하는 AI 에이전트\n실제 사람처럼 행동하는 NPC"),
    ]

    for i, (icon, title, desc) in enumerate(cards):
        cx = 0.5 + i * 4.27
        add_rect(s, cx, 1.5, 3.9, 4.5, BG_CARD)
        add_rect(s, cx, 1.5, 3.9, 0.06, C_RED)
        add_text(s, icon, cx + 0.2, 1.7, 0.8, 0.7, font_size=32)
        add_text(s, title, cx + 0.2, 2.35, 3.5, 0.6,
                 font_size=20, bold=True, color=C_WHITE, font_name="Malgun Gothic")
        add_divider(s, cx + 0.2, 3.0, 3.4, C_RED)
        add_text(s, desc, cx + 0.2, 3.15, 3.5, 2.5,
                 font_size=16, color=C_DIMWHITE, font_name="Malgun Gothic")

    # Bottom conclusion
    add_rect(s, 0.5, 6.2, 12.33, 0.9, RGBColor(0x1C, 0x10, 0x10))
    add_rect(s, 0.5, 6.2, 0.06, 0.9, C_RED)
    add_text(s,
        '→  "LLM을 활용해 진짜 살아있는 NPC가 존재하는 내러티브 게임 엔진을 만들자!"',
        0.8, 6.32, 12, 0.6,
        font_size=19, bold=True, color=C_GOLD, font_name="Malgun Gothic")
    return s


def slide_concept(prs):
    """Slide 4: Core Concept"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "핵심 컨셉", 0.6, 0.25, 8, 0.7,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, '"인형의 집"에서 탈출하라', 0.6, 0.72, 8, 0.4,
             font_size=16, color=C_RED, italic=True, font_name="Malgun Gothic")

    # Left: Game concept
    add_rect(s, 0.5, 1.4, 5.8, 5.7, BG_CARD)
    add_rect(s, 0.5, 1.4, 5.8, 0.06, C_RED)
    add_text(s, "🎮  게임 소개", 0.8, 1.55, 5.2, 0.5,
             font_size=20, bold=True, color=C_WHITE, font_name="Malgun Gothic")

    game_points = [
        "플레이어는 불가사의한 인형의 집에 갇힘",
        "버튼 눈을 가진 인형 가족들이 당신을 감시",
        "50턴 안에 탈출하거나 인형이 됨",
        "NPC와 대화·행동으로 단서 수집",
        "5가지 엔딩 중 하나로 귀결",
    ]
    for i, pt in enumerate(game_points):
        add_rect(s, 0.7, 2.25 + i * 0.87, 0.06, 0.5, C_RED)
        add_text(s, pt, 0.9, 2.2 + i * 0.87, 5.2, 0.55,
                 font_size=17, color=C_DIMWHITE, font_name="Malgun Gothic")

    # Right: AI Core
    add_rect(s, 6.7, 1.4, 6.1, 5.7, BG_CARD)
    add_rect(s, 6.7, 1.4, 6.1, 0.06, C_TEAL)
    add_text(s, "🧠  AI 핵심 메커니즘", 7.0, 1.55, 5.5, 0.5,
             font_size=20, bold=True, color=C_WHITE, font_name="Malgun Gothic")

    ai_items = [
        ("기억 (Memory)", "모든 대화·사건을 NPC가 기억"),
        ("반성 (Reflection)", "하루 끝 NPC가 스스로 반성·학습"),
        ("계획 (Planning)", "다음날 행동을 미리 계획"),
        ("자율 대화 (Dialogue)", "NPC들끼리 밤에 자체 대화"),
        ("도구 호출 (Tools)", "LLM이 interact/action/use 선택"),
    ]
    for i, (label, desc) in enumerate(ai_items):
        add_rect(s, 6.9, 2.25 + i * 0.87, 0.06, 0.5, C_TEAL)
        add_text(s, label, 7.1, 2.2 + i * 0.87, 2.5, 0.4,
                 font_size=16, bold=True, color=C_TEAL, font_name="Malgun Gothic")
        add_text(s, desc, 9.6, 2.2 + i * 0.87, 3.0, 0.4,
                 font_size=15, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_tech_stack(prs):
    """Slide 5: Tech Stack"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "기술 스택", 0.6, 0.25, 8, 0.7,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "What we built with", 0.6, 0.72, 5, 0.4,
             font_size=16, color=C_TEAL, font_name="Malgun Gothic")

    categories = [
        ("🔥  AI / LLM", C_RED, [
            "OpenAI GPT-4", "Generative Agents Framework",
            "Tool Calling (Function Calling)", "Memory Streams",
        ]),
        ("⚡  Backend", C_TEAL, [
            "FastAPI 0.100+", "SQLAlchemy 2.0 (ORM)",
            "Alembic (DB Migration)", "Uvicorn (ASGI)",
        ]),
        ("🗄️  Data & Infra", C_GOLD, [
            "PostgreSQL", "Redis (Cache/State)",
            "Docker Compose", "PyYAML (Scenario DSL)",
        ]),
        ("🔧  Dev Tools", C_PURPLE, [
            "Pydantic v2 (Schema)", "python-dotenv",
            "Ngrok (Tunneling)", "Jupyter Notebook (Testing)",
        ]),
    ]

    for i, (cat_title, accent, items) in enumerate(categories):
        col = i % 2
        row = i // 2
        cx = 0.5 + col * 6.4
        cy = 1.4 + row * 2.8

        add_rect(s, cx, cy, 6.0, 2.55, BG_CARD)
        add_rect(s, cx, cy, 6.0, 0.06, accent)
        add_text(s, cat_title, cx + 0.2, cy + 0.12, 5.5, 0.45,
                 font_size=18, bold=True, color=accent, font_name="Malgun Gothic")
        add_divider(s, cx + 0.2, cy + 0.62, 5.6, accent)
        for j, item in enumerate(items):
            add_text(s, f"▸  {item}", cx + 0.25, cy + 0.75 + j * 0.4, 5.5, 0.38,
                     font_size=15, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_architecture(prs):
    """Slide 6: System Architecture"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "시스템 아키텍처", 0.6, 0.25, 8, 0.7,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "8-Layer Architecture", 0.6, 0.72, 5, 0.4,
             font_size=16, color=C_TEAL, font_name="Malgun Gothic")

    layers = [
        ("API Layer", "FastAPI REST Endpoints", C_TEAL),
        ("Service Layer", "GameService / ScenarioLoader", C_TEAL),
        ("Controller", "DayController / NightController", C_GOLD),
        ("Tools", "interact / action / use", C_GOLD),
        ("LLM Engine", "Prompt Builder → GPT-4 → Parser", C_RED),
        ("Rule Engine", "MemoryRules / LockManager / ItemAcquirer", C_RED),
        ("State Manager", "WorldStateManager / StatusEffectManager", C_PURPLE),
        ("Narrative", "NarrativeLayer → Final Prose", C_GREEN),
    ]

    for i, (name, desc, color) in enumerate(layers):
        col = i % 2
        row = i // 2
        cx = 0.5 + col * 6.4
        cy = 1.4 + row * 1.45

        add_rect(s, cx, cy, 5.9, 1.25, BG_CARD)
        add_rect(s, cx, cy, 0.08, 1.25, color)
        add_text(s, name, cx + 0.25, cy + 0.1, 5.4, 0.5,
                 font_size=20, bold=True, color=color, font_name="Malgun Gothic")
        add_text(s, desc, cx + 0.25, cy + 0.6, 5.4, 0.5,
                 font_size=14, color=C_DIMWHITE, font_name="Malgun Gothic")

        # Arrow down within same column
        if i < 6 and row < 3:
            ax = cx + 2.9
            ay = cy + 1.25
            add_rect(s, ax, ay, 0.04, 0.2, color)
    return s


def slide_generative_agents(prs):
    """Slide 7: Generative Agents"""
    s = blank_slide(prs)
    fill_bg(s, BG_NIGHT)

    add_rect(s, 0, 0, 13.33, 1.2, RGBColor(0x08, 0x12, 0x1E))
    add_text(s, "Generative Agents", 0.6, 0.22, 10, 0.75,
             font_size=40, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "NPC가 진짜 인간처럼 기억하고, 반성하고, 계획한다  (Park et al. 2023 inspired)", 0.6, 0.72, 12, 0.4,
             font_size=15, color=C_TEAL, italic=True, font_name="Malgun Gothic")

    # Four agent phases
    phases = [
        ("💭", "Memory Stream", C_TEAL,
         "모든 대화와 사건이\n메모리 스트림에 저장\n시간·중요도·유사도 기반 검색"),
        ("🔮", "Reflection", C_PURPLE,
         "하루 끝에 NPC가 스스로\n중요한 기억을 반성·통합\n고차원 인사이트 생성"),
        ("📋", "Planning", C_GOLD,
         "반성을 바탕으로\n다음날 행동 계획 수립\n장기 목표 + 단기 행동"),
        ("💬", "Group Dialogue", C_RED,
         "밤마다 NPC들이 자체 대화\n6번의 자율 발화\n관계·의심도 자동 변화"),
    ]

    for i, (icon, title, color, desc) in enumerate(phases):
        cx = 0.5 + i * 3.2
        add_rect(s, cx, 1.5, 2.95, 4.8, RGBColor(0x0F, 0x1C, 0x30))
        add_rect(s, cx, 1.5, 2.95, 0.06, color)
        add_text(s, icon, cx + 0.2, 1.65, 0.7, 0.7, font_size=28)
        add_text(s, title, cx + 0.15, 2.3, 2.65, 0.5,
                 font_size=20, bold=True, color=color, font_name="Malgun Gothic")
        add_divider(s, cx + 0.15, 2.85, 2.65, color)
        add_text(s, desc, cx + 0.15, 3.0, 2.65, 3.0,
                 font_size=15, color=C_DIMWHITE, font_name="Malgun Gothic")

        if i < 3:
            add_text(s, "→", cx + 3.0, 3.5, 0.3, 0.4,
                     font_size=22, bold=True, color=C_SLATE, align=PP_ALIGN.CENTER)

    # Bottom note
    add_rect(s, 0.5, 6.5, 12.33, 0.65, RGBColor(0x0A, 0x18, 0x2C))
    add_rect(s, 0.5, 6.5, 0.06, 0.65, C_TEAL)
    add_text(s,
        "각 NPC (Eleanor, Arthur, Lucas, Margaret, Baron)마다 독립적인 메모리 스트림 운영 → 완전히 다른 인격과 반응",
        0.8, 6.55, 12, 0.55,
        font_size=16, color=C_DIMWHITE, italic=True, font_name="Malgun Gothic")
    return s


def slide_day_pipeline(prs):
    """Slide 8: Day Pipeline"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "Day Pipeline", 0.6, 0.22, 8, 0.75,
             font_size=40, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "플레이어 턴 처리 흐름", 0.6, 0.72, 6, 0.4,
             font_size=16, color=C_GOLD, font_name="Malgun Gothic")

    steps = [
        ("1", "User Input", "자연어 입력 접수", C_TEAL),
        ("2", "Tool Calling", "LLM이 도구 선택\ninteract / action / use", C_TEAL),
        ("3", "LLM Execution", "선택된 도구로\nGPT-4 프롬프트 실행", C_GOLD),
        ("4", "Delta Parse", "JSON 응답에서\nStateDelta 추출", C_GOLD),
        ("5", "Rule Engine", "memory_rules.yaml\n조건부 효과 자동 적용", C_RED),
        ("6", "State Apply", "WorldStateManager\n게임 상태 업데이트", C_RED),
        ("7", "Narrative", "NarrativeLayer\n산문 형태 최종 출력", C_GREEN),
    ]

    # Flow line
    add_rect(s, 0.5, 4.0, 12.33, 0.04, C_SLATE)

    for i, (num, title, desc, color) in enumerate(steps):
        cx = 0.45 + i * 1.77
        # Circle number
        add_rect(s, cx, 3.6, 0.7, 0.7, color)
        add_text(s, num, cx, 3.62, 0.7, 0.65,
                 font_size=22, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
        # Title
        add_text(s, title, cx - 0.1, 4.25, 1.8, 0.5,
                 font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
        # Desc
        add_text(s, desc, cx - 0.1, 4.75, 1.8, 0.8,
                 font_size=11, color=C_DIMWHITE, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")

    # Intent system callout
    add_rect(s, 0.5, 1.35, 5.8, 2.3, BG_CARD)
    add_rect(s, 0.5, 1.35, 5.8, 0.06, C_GOLD)
    add_text(s, "🎯  Intent Classification", 0.75, 1.48, 5.3, 0.5,
             font_size=18, bold=True, color=C_GOLD, font_name="Malgun Gothic")
    intents = [
        ("investigate", C_RED, "의심 행동 → 서스펜션 +5"),
        ("obey", C_GREEN, "순응 → NPC 호감 +, 인간성 -"),
        ("rebel", C_RED, "반항 → NPC 호감 -, 의심 +"),
        ("reveal", C_TEAL, "진실 탐구 → NPC humanity +"),
        ("neutral", C_SLATE, "일상 대화 → 변화 없음"),
    ]
    for j, (intent, ic, effect) in enumerate(intents):
        add_rect(s, 0.7, 2.05 + j * 0.3, 0.06, 0.22, ic)
        add_text(s, f"{intent}  →  {effect}", 0.9, 2.05 + j * 0.3, 5.1, 0.28,
                 font_size=13, color=C_DIMWHITE, font_name="Malgun Gothic")

    # Tool callout
    add_rect(s, 7.0, 1.35, 5.8, 2.3, BG_CARD)
    add_rect(s, 7.0, 1.35, 5.8, 0.06, C_TEAL)
    add_text(s, "🔧  Tool System", 7.25, 1.48, 5.3, 0.5,
             font_size=18, bold=True, color=C_TEAL, font_name="Malgun Gothic")
    tools_list = [
        ("interact(target, type)", "NPC 대화 또는 물체 조사"),
        ("action(action)", "비상호 행동 수행"),
        ("use(item, action, target)", "아이템 사용 (acquire/usage)"),
    ]
    for j, (tool, tdesc) in enumerate(tools_list):
        add_rect(s, 7.2, 2.05 + j * 0.54, 0.06, 0.4, C_TEAL)
        add_text(s, tool, 7.4, 2.05 + j * 0.54, 5.1, 0.3,
                 font_size=13, bold=True, color=C_TEAL, font_name="Malgun Gothic")
        add_text(s, tdesc, 7.4, 2.37 + j * 0.54, 5.1, 0.25,
                 font_size=12, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_night_pipeline(prs):
    """Slide 9: Night Pipeline"""
    s = blank_slide(prs)
    fill_bg(s, BG_NIGHT)

    add_rect(s, 0, 0, 13.33, 1.2, RGBColor(0x06, 0x0E, 0x1C))
    add_text(s, "Night Pipeline", 0.6, 0.22, 8, 0.75,
             font_size=40, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "NPC 자율 행동 — 밤은 AI가 지배한다", 0.6, 0.72, 8, 0.4,
             font_size=16, color=C_PURPLE, italic=True, font_name="Malgun Gothic")

    # Phase cards - horizontal
    phases = [
        ("Phase 1", "Reflection\n반성", C_PURPLE,
         "각 NPC가 오늘 하루를 돌아봄\n중요한 기억·감정 정리\n인사이트 추출"),
        ("Phase 2", "Planning\n계획", C_GOLD,
         "반성 내용 기반\n내일 행동 계획 수립\n동적 감정 포맷 적용"),
        ("Phase 3", "Group Dialogue\n집단 대화", C_RED,
         "NPC 2명이 무작위 대화\n6번 발화 (라운드로빈)\n플레이어 관련 음모 가능"),
        ("Phase 4", "Impact Analysis\n영향 분석", C_TEAL,
         "대화에서 스탯 변화 추출\nparse_stat_changes_text()\nnight_delta 누적 적용"),
    ]

    for i, (phase_num, title, color, desc) in enumerate(phases):
        cx = 0.4 + i * 3.25
        add_rect(s, cx, 1.45, 3.0, 4.9, RGBColor(0x0C, 0x18, 0x28))
        add_rect(s, cx, 1.45, 3.0, 0.06, color)
        add_text(s, phase_num, cx + 0.15, 1.58, 2.7, 0.38,
                 font_size=13, bold=True, color=color, font_name="Malgun Gothic")
        add_text(s, title, cx + 0.15, 1.98, 2.7, 0.8,
                 font_size=21, bold=True, color=C_WHITE, font_name="Malgun Gothic")
        add_divider(s, cx + 0.15, 2.82, 2.7, color)
        add_text(s, desc, cx + 0.15, 3.0, 2.7, 3.0,
                 font_size=15, color=C_DIMWHITE, font_name="Malgun Gothic")

        if i < 3:
            add_text(s, "▶", 3.44 + i * 3.25, 3.8, 0.4, 0.45,
                     font_size=18, color=C_SLATE, align=PP_ALIGN.CENTER)

    # Bottom: Night's impact
    add_rect(s, 0.4, 6.55, 12.53, 0.65, RGBColor(0x10, 0x06, 0x20))
    add_rect(s, 0.4, 6.55, 0.06, 0.65, C_PURPLE)
    add_text(s,
        "플레이어가 잠든 사이 NPC들은 스스로 생각하고, 계획하고, 서로 음모를 꾸민다 — 매 밤 게임이 진화한다",
        0.7, 6.6, 12, 0.55,
        font_size=16, color=C_DIMWHITE, italic=True, font_name="Malgun Gothic")
    return s


def slide_dynamic_stats(prs):
    """Slide 10: Dynamic Stats & Rule Engine"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "동적 스탯 & 규칙 엔진", 0.6, 0.22, 10, 0.75,
             font_size=40, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "Scenario YAML만 수정하면 코드 변경 없이 새 게임 시스템 구축", 0.6, 0.72, 12, 0.4,
             font_size=16, color=C_GREEN, font_name="Malgun Gothic")

    # Left: Before vs After
    add_rect(s, 0.5, 1.4, 5.8, 5.7, BG_CARD)
    add_rect(s, 0.5, 1.4, 5.8, 0.06, C_SLATE)
    add_text(s, "🔴  Before (v2) — 하드코딩", 0.75, 1.55, 5.3, 0.5,
             font_size=17, bold=True, color=C_RED, font_name="Malgun Gothic")

    before_code = "# 코드에 고정된 스탯\nnpc.trust = 50\nnpc.suspicion = 0\nnpc.fear = 30\nnpc.humanity = 100\n\n# 스탯 추가 시 코드 전체 수정"
    add_rect(s, 0.7, 2.1, 5.3, 2.1, RGBColor(0x1A, 0x08, 0x08))
    add_text(s, before_code, 0.9, 2.15, 5.0, 2.0,
             font_size=13, color=RGBColor(0xFF, 0x88, 0x66), font_name="Courier New")

    add_text(s, "문제점:", 0.75, 4.3, 5.3, 0.4,
             font_size=15, bold=True, color=C_RED, font_name="Malgun Gothic")
    problems = ["새 시나리오마다 코드 변경 필수", "스탯 추가·제거 불가", "범용성 zero"]
    for i, p in enumerate(problems):
        add_text(s, f"  ✗  {p}", 0.75, 4.7 + i * 0.35, 5.3, 0.33,
                 font_size=14, color=C_DIMWHITE, font_name="Malgun Gothic")

    # Right: After
    add_rect(s, 6.8, 1.4, 6.0, 5.7, BG_CARD)
    add_rect(s, 6.8, 1.4, 6.0, 0.06, C_GREEN)
    add_text(s, "🟢  After (v3) — YAML 동적 정의", 7.05, 1.55, 5.5, 0.5,
             font_size=17, bold=True, color=C_GREEN, font_name="Malgun Gothic")

    after_code = "# npcs.yaml\ninitial_stats:\n  affection: 50\n  fear: 80\n  humanity: 0\n\n# 어떤 스탯도 추가 가능!"
    add_rect(s, 7.0, 2.1, 5.5, 2.1, RGBColor(0x08, 0x18, 0x10))
    add_text(s, after_code, 7.2, 2.15, 5.2, 2.0,
             font_size=13, color=RGBColor(0x66, 0xFF, 0x88), font_name="Courier New")

    add_text(s, "장점:", 7.05, 4.3, 5.5, 0.4,
             font_size=15, bold=True, color=C_GREEN, font_name="Malgun Gothic")
    benefits = ["시나리오 제작자가 임의 스탯 정의", "코드 변경 없이 새 게임 시스템", "완전한 범용 엔진"]
    for i, b in enumerate(benefits):
        add_text(s, f"  ✓  {b}", 7.05, 4.7 + i * 0.35, 5.5, 0.33,
                 font_size=14, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_rule_engine(prs):
    """Slide 11: Rule Engine Detail"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "규칙 엔진 (Memory Rules)", 0.6, 0.22, 10, 0.75,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "플레이어 행동 의도에 따라 자동으로 게임 상태 변화", 0.6, 0.72, 12, 0.4,
             font_size=16, color=C_GOLD, font_name="Malgun Gothic")

    # Left side: YAML example
    add_rect(s, 0.5, 1.4, 6.0, 5.7, BG_CARD)
    add_rect(s, 0.5, 1.4, 6.0, 0.06, C_GOLD)
    add_text(s, "📄  memory_rules.yaml", 0.75, 1.55, 5.5, 0.45,
             font_size=17, bold=True, color=C_GOLD, font_name="Malgun Gothic")

    yaml_code = """- rule_id: investigate_increase_suspicion
  when: "intent == 'investigate'"
  effects:
    - type: var_add
      key: vars.suspicion_level
      value: 5
    - type: var_add
      key: vars.humanity
      value: 2

- rule_id: obey_gain_affection
  when: "intent == 'obey'"
  effects:
    - type: npc_stat_add
      target: stepmother
      key: affection
      value: 10
    - type: var_add
      key: vars.humanity
      value: -5"""
    add_rect(s, 0.7, 2.05, 5.6, 4.85, RGBColor(0x10, 0x10, 0x18))
    add_text(s, yaml_code, 0.85, 2.1, 5.4, 4.75,
             font_size=12, color=RGBColor(0x88, 0xCC, 0xFF), font_name="Courier New")

    # Right side: effect types
    add_rect(s, 6.9, 1.4, 5.9, 5.7, BG_CARD)
    add_rect(s, 6.9, 1.4, 5.9, 0.06, C_GOLD)
    add_text(s, "⚙️  Effect Types", 7.15, 1.55, 5.4, 0.45,
             font_size=17, bold=True, color=C_GOLD, font_name="Malgun Gothic")

    effect_types = [
        ("var_add", "world vars 수치 증감", C_TEAL),
        ("var_set", "world vars 직접 설정", C_TEAL),
        ("npc_stat_add", "NPC 개별 스탯 증감", C_RED),
        ("npc_status_set", "NPC 상태 변경 (sleeping...)", C_RED),
        ("flag_set", "게임 플래그 토글", C_GOLD),
        ("add_item", "인벤토리에 아이템 추가", C_GREEN),
        ("remove_item", "인벤토리에서 제거", C_GREEN),
        ("memory_add", "NPC 메모리에 사건 삽입", C_PURPLE),
    ]

    for i, (etype, edesc, ecolor) in enumerate(effect_types):
        add_rect(s, 7.1, 2.1 + i * 0.62, 0.06, 0.45, ecolor)
        add_text(s, etype, 7.3, 2.1 + i * 0.62, 2.5, 0.3,
                 font_size=13, bold=True, color=ecolor, font_name="Courier New")
        add_text(s, edesc, 7.3, 2.38 + i * 0.62, 5.2, 0.28,
                 font_size=12, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_characters(prs):
    """Slide 12: NPC Characters"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "시나리오: 코렐라인 — 등장인물", 0.6, 0.22, 10, 0.75,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, '"인형의 집" 5인의 섬뜩한 가족들', 0.6, 0.72, 8, 0.4,
             font_size=16, color=C_RED, italic=True, font_name="Malgun Gothic")

    characters = [
        ("👁️", "새엄마\nEleanor", C_RED,
         "역할: 주 적대자\n인형 가족 완성을 목표\n3단계 행동 변화\n(유혹→시험→처벌)"),
        ("🪨", "새아빠\nArthur", RGBColor(0x6C, 0x8E, 0xAD),
         "역할: 물리적 집행자\n침묵·순종 가치관\n단서 힌트 제공\n(신문·과거 물건)"),
        ("🪆", "남동생\nLucas", C_TEAL,
         "역할: 반인반인형\n인간성·인형성 갈등\n동반자 혹은 희생자\n탈출 정보 보유"),
        ("📜", "할머니\nMargaret", C_GOLD,
         "역할: 옛 지배자\n집 비밀 보유\n불꽃 약점 정보\n억압된 동맹군"),
        ("🐕", "바론\n(개)", C_GREEN,
         "역할: 유일한 순수\n진짜 인간 감지\n보호 힌트 제공\n아이템: 간식"),
    ]

    for i, (icon, name, color, desc) in enumerate(characters):
        cx = 0.4 + i * 2.55
        add_rect(s, cx, 1.4, 2.3, 5.7, BG_CARD)
        add_rect(s, cx, 1.4, 2.3, 0.06, color)
        add_text(s, icon, cx + 0.75, 1.55, 0.9, 0.7, font_size=30, align=PP_ALIGN.CENTER)
        add_text(s, name, cx + 0.1, 2.25, 2.1, 0.7,
                 font_size=17, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
        add_divider(s, cx + 0.1, 3.0, 2.1, color)
        add_text(s, desc, cx + 0.15, 3.15, 2.05, 3.6,
                 font_size=13, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_world_vars(prs):
    """Slide 13: World Variables & Mechanics"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "핵심 게임 메커니즘", 0.6, 0.22, 10, 0.75,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "두 개의 핵심 카운터가 운명을 결정한다", 0.6, 0.72, 10, 0.4,
             font_size=16, color=C_RED, italic=True, font_name="Malgun Gothic")

    # Humanity gauge (left)
    add_rect(s, 0.5, 1.45, 5.8, 5.6, BG_CARD)
    add_rect(s, 0.5, 1.45, 5.8, 0.06, C_TEAL)
    add_text(s, "인간성 (Humanity)", 0.75, 1.62, 5.3, 0.5,
             font_size=22, bold=True, color=C_TEAL, font_name="Malgun Gothic")
    add_text(s, "0 ~ 100", 0.75, 2.12, 2, 0.38,
             font_size=16, color=C_DIMWHITE, font_name="Malgun Gothic")

    # Bar visualization
    add_rect(s, 0.75, 2.6, 5.2, 0.45, RGBColor(0x20, 0x20, 0x30))
    add_rect(s, 0.75, 2.6, 3.0, 0.45, C_TEAL)  # 60% human
    add_text(s, "60%", 3.9, 2.63, 1.0, 0.38,
             font_size=16, bold=True, color=C_TEAL, font_name="Malgun Gothic")
    add_text(s, "100 = 완전한 인간\n0 = 인형이 됨 → 게임 오버", 0.75, 3.22, 5.2, 0.7,
             font_size=15, color=C_DIMWHITE, font_name="Malgun Gothic")

    humanity_effects = [
        ("순응·복종", "↓", C_RED),
        ("허브 냄새 맡기", "↑", C_GREEN),
        ("인형 물건 접촉", "↓", C_RED),
        ("가족 사진 발견", "↑", C_GREEN),
    ]
    for i, (act, arrow, color) in enumerate(humanity_effects):
        add_rect(s, 0.75, 4.1 + i * 0.5, 0.06, 0.35, color)
        add_text(s, f"{act}  {arrow}", 0.95, 4.12 + i * 0.5, 4.8, 0.33,
                 font_size=14, color=C_DIMWHITE, font_name="Malgun Gothic")

    # Suspicion gauge (right)
    add_rect(s, 6.9, 1.45, 5.9, 5.6, BG_CARD)
    add_rect(s, 6.9, 1.45, 5.9, 0.06, C_RED)
    add_text(s, "의심도 (Suspicion)", 7.15, 1.62, 5.4, 0.5,
             font_size=22, bold=True, color=C_RED, font_name="Malgun Gothic")
    add_text(s, "0 ~ 100", 7.15, 2.12, 2, 0.38,
             font_size=16, color=C_DIMWHITE, font_name="Malgun Gothic")

    add_rect(s, 7.15, 2.6, 5.2, 0.45, RGBColor(0x20, 0x20, 0x30))
    add_rect(s, 7.15, 2.6, 1.6, 0.45, C_RED)  # 30% suspicion
    add_text(s, "30%", 8.9, 2.63, 1.0, 0.38,
             font_size=16, bold=True, color=C_RED, font_name="Malgun Gothic")
    add_text(s, "0 = 완전한 신뢰\n100 = 붙잡힘 → 게임 오버", 7.15, 3.22, 5.2, 0.7,
             font_size=15, color=C_DIMWHITE, font_name="Malgun Gothic")

    susp_effects = [
        ("조사 행동", "↑", C_RED),
        ("순종적 행동", "↓", C_GREEN),
        ("탈출 시도", "↑↑", C_RED),
        ("식사 칭찬", "↓", C_GREEN),
    ]
    for i, (act, arrow, color) in enumerate(susp_effects):
        add_rect(s, 7.15, 4.1 + i * 0.5, 0.06, 0.35, color)
        add_text(s, f"{act}  {arrow}", 7.35, 4.12 + i * 0.5, 4.8, 0.33,
                 font_size=14, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_endings(prs):
    """Slide 14: Ending System"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "엔딩 시스템 — 5가지 결말", 0.6, 0.22, 10, 0.75,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "당신의 선택이 운명을 결정한다", 0.6, 0.72, 8, 0.4,
             font_size=16, color=C_GOLD, italic=True, font_name="Malgun Gothic")

    endings = [
        ("stealth_exit", "🤫  은밀한 탈출", C_GREEN,
         "✓ VICTORY",
         "진정제로 가족 재우기\n비밀 열쇠 획득\n완벽한 위장 탈출"),
        ("chaotic_breakout", "🔥  혼돈의 탈출", C_ORANGE,
         "✓ VICTORY",
         "기름 + 라이터로 새엄마에게\n불 지르기\n화염 속 도주"),
        ("sibling_sacrifice", "💔  동생의 희생", C_TEAL,
         "✓ VICTORY",
         "남동생 루카스의\n자기희생으로 탈출\n쌍방 신뢰 필요"),
        ("unfinished_doll", "🪆  미완성 인형", C_RED,
         "✗ FAILURE",
         "인간성 0 도달\n완전한 인형이 됨\n영원히 집에 남음"),
        ("caught_confined", "⛓️  포획·감금", C_PURPLE,
         "✗ FAILURE",
         "의심도 100 달성\n가족에게 붙잡힘\n감금 종료"),
    ]

    for i, (eid, title, color, result, desc) in enumerate(endings):
        col = i % 3
        row = i // 3
        cx = 0.4 + col * 4.3
        cy = 1.4 + row * 3.1

        add_rect(s, cx, cy, 4.0, 2.85, BG_CARD)
        add_rect(s, cx, cy, 4.0, 0.06, color)
        add_text(s, title, cx + 0.15, cy + 0.15, 3.7, 0.5,
                 font_size=18, bold=True, color=color, font_name="Malgun Gothic")
        # Result badge
        badge_color = C_GREEN if "✓" in result else C_RED
        add_rect(s, cx + 0.15, cy + 0.68, 1.5, 0.32, badge_color)
        add_text(s, result, cx + 0.15, cy + 0.68, 1.5, 0.32,
                 font_size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
        add_text(s, desc, cx + 0.15, cy + 1.12, 3.7, 1.65,
                 font_size=14, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_dev_status(prs):
    """Slide 15: Development Status"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "개발 현황", 0.6, 0.22, 8, 0.75,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "Recent Commits & Architecture Evolution", 0.6, 0.72, 8, 0.4,
             font_size=16, color=C_TEAL, font_name="Malgun Gothic")

    # Recent commits timeline
    add_rect(s, 0.5, 1.4, 7.5, 5.7, BG_CARD)
    add_rect(s, 0.5, 1.4, 7.5, 0.06, C_TEAL)
    add_text(s, "📋  최근 커밋 히스토리", 0.75, 1.55, 7.0, 0.45,
             font_size=17, bold=True, color=C_TEAL, font_name="Malgun Gothic")

    commits = [
        ("eee7735", "coraline_v3 YAML 스키마 파이프라인 호환성 수정", "refactor", C_GOLD),
        ("9f6d788", "house_on_fire → flags, status_effects → vars 이동", "refactor", C_GOLD),
        ("6a6cb39", "StatusEffectManager status 전용 재설계", "refactor", C_TEAL),
        ("605cb08", "StatusEffectManager & ItemAcquirer 서비스에 추가", "feat", C_GREEN),
        ("d8136f4", "dev/raphael 브랜치 병합", "merge", C_PURPLE),
    ]

    for i, (sha, msg, tag, color) in enumerate(commits):
        cy = 2.1 + i * 0.95
        add_rect(s, 0.7, cy, 0.06, 0.7, color)
        # Tag badge
        add_rect(s, 0.9, cy + 0.1, 1.0, 0.28, color)
        add_text(s, tag, 0.9, cy + 0.1, 1.0, 0.28,
                 font_size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
        add_text(s, sha, 2.05, cy + 0.1, 1.0, 0.3,
                 font_size=12, color=C_SLATE, font_name="Courier New")
        add_text(s, msg, 0.9, cy + 0.42, 6.8, 0.35,
                 font_size=14, color=C_DIMWHITE, font_name="Malgun Gothic")

    # Right: Key metrics
    add_rect(s, 8.3, 1.4, 4.6, 5.7, BG_CARD)
    add_rect(s, 8.3, 1.4, 4.6, 0.06, C_GOLD)
    add_text(s, "📊  프로젝트 규모", 8.55, 1.55, 4.1, 0.45,
             font_size=17, bold=True, color=C_GOLD, font_name="Malgun Gothic")

    metrics = [
        ("~5,855", "Python 코드 라인"),
        ("8", "아키텍처 레이어"),
        ("5", "NPC 캐릭터"),
        ("5", "게임 엔딩"),
        ("12+", "인게임 아이템"),
        ("50", "최대 플레이 턴"),
        ("10", "문서/다이어그램"),
    ]

    for i, (num, label) in enumerate(metrics):
        add_rect(s, 8.5, 2.12 + i * 0.72, 4.2, 0.6, RGBColor(0x1A, 0x1A, 0x2A))
        add_text(s, num, 8.65, 2.16 + i * 0.72, 1.5, 0.5,
                 font_size=26, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
        add_text(s, label, 10.35, 2.24 + i * 0.72, 2.2, 0.35,
                 font_size=13, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_future(prs):
    """Slide 16: Future Plans"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    add_rect(s, 0, 0, 13.33, 1.2, BG_CARD)
    add_text(s, "향후 계획", 0.6, 0.22, 8, 0.75,
             font_size=38, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    add_text(s, "Next Steps for DeusExMachina", 0.6, 0.72, 8, 0.4,
             font_size=16, color=C_TEAL, font_name="Malgun Gothic")

    roadmap = [
        ("🎮  새 시나리오 확장", C_TEAL,
         ["YAML DSL 기반 신규 시나리오 제작", "다양한 장르 (미스터리, SF, 판타지)", "멀티 플레이어 지원 검토"]),
        ("🧠  AI 고도화", C_PURPLE,
         ["더 정교한 메모리 검색 (임베딩)", "NPC 감정 모델 강화", "장기 계획 개선"]),
        ("⚙️  성능 최적화", C_GOLD,
         ["LLM 호출 비용 절감 (캐싱)", "응답 속도 개선", "병렬 NPC 처리"]),
        ("📱  사용자 경험", C_GREEN,
         ["웹 프론트엔드 구축", "실시간 상태 시각화", "저장/불러오기 UI"]),
    ]

    for i, (title, color, items) in enumerate(roadmap):
        col = i % 2
        row = i // 2
        cx = 0.5 + col * 6.4
        cy = 1.4 + row * 2.9

        add_rect(s, cx, cy, 6.0, 2.65, BG_CARD)
        add_rect(s, cx, cy, 6.0, 0.06, color)
        add_text(s, title, cx + 0.2, cy + 0.15, 5.6, 0.5,
                 font_size=20, bold=True, color=color, font_name="Malgun Gothic")
        add_divider(s, cx + 0.2, cy + 0.7, 5.6, color)
        for j, item in enumerate(items):
            add_rect(s, cx + 0.2, cy + 0.85 + j * 0.52, 0.06, 0.35, color)
            add_text(s, item, cx + 0.4, cy + 0.83 + j * 0.52, 5.4, 0.38,
                     font_size=15, color=C_DIMWHITE, font_name="Malgun Gothic")
    return s


def slide_closing(prs):
    """Slide 17: Thank You / Closing"""
    s = blank_slide(prs)
    fill_bg(s, BG_DARK)

    # Left accent
    add_rect(s, 0, 0, 0.5, 7.5, C_RED)

    # Grid lines
    for i in range(1, 8):
        add_rect(s, 0.5, i * 0.95, 12.83, 0.005, RGBColor(0x20, 0x20, 0x30))

    # Main message
    add_text(s, "모든 선택이 이야기를 만든다", 1.0, 1.5, 11, 1.2,
             font_size=56, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")

    add_divider(s, 2.5, 3.0, 8.33, C_RED, 0.04)

    add_text(s, "DeusExMachina", 1.0, 3.2, 11, 0.9,
             font_size=42, bold=True, color=C_RED, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")

    add_text(s,
        "AI가 살아 숨쉬는 NPC · 당신만의 이야기 · 50턴의 운명",
        1.0, 4.15, 11, 0.6,
        font_size=20, color=C_PARCHMENT, align=PP_ALIGN.CENTER, italic=True, font_name="Malgun Gothic")

    add_divider(s, 2.5, 5.0, 8.33, C_SLATE)

    add_text(s,
        "FastAPI  ·  OpenAI GPT-4  ·  Generative Agents  ·  PostgreSQL  ·  Redis",
        1.0, 5.2, 11, 0.5,
        font_size=15, color=C_SLATE, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")

    add_text(s, "LIKELION × DeusExMachina Team  ·  2025",
             1.0, 6.3, 11, 0.45,
             font_size=14, color=C_SLATE, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")

    add_rect(s, 12.33, 6.5, 1.0, 0.04, C_GOLD)
    add_text(s, "인형의 집에서 탈출하라", 9.3, 6.6, 3.7, 0.5,
             font_size=13, color=C_GOLD, italic=True,
             align=PP_ALIGN.RIGHT, font_name="Malgun Gothic")
    return s


# ─── MAIN ─────────────────────────────────────────────────────────────────────

def build_presentation():
    prs = new_prs()

    print("Building slides...")

    # 1. Title
    slide_title(prs)
    print("  [1/17] Title")

    # 2. Section: Background
    slide_section_divider(prs, 1, "프로젝트 배경", "왜 이걸 만들었나?")
    print("  [2/17] Section divider")

    # 3. Background
    slide_project_background(prs)
    print("  [3/17] Background")

    # 4. Section: Concept
    slide_section_divider(prs, 2, "핵심 컨셉", '"코렐라인" — 인형의 집 탈출 어드벤처')
    print("  [4/17] Section divider")

    # 5. Concept
    slide_concept(prs)
    print("  [5/17] Concept")

    # 6. Section: Tech
    slide_section_divider(prs, 3, "기술 스택 & 아키텍처", "AI 내러티브 엔진의 구조")
    print("  [6/17] Section divider")

    # 7. Tech Stack
    slide_tech_stack(prs)
    print("  [7/17] Tech Stack")

    # 8. Architecture
    slide_architecture(prs)
    print("  [8/17] Architecture")

    # 9. Section: Core Features
    slide_section_divider(prs, 4, "핵심 기능", "Generative Agents & 파이프라인")
    print("  [9/17] Section divider")

    # 10. Generative Agents
    slide_generative_agents(prs)
    print("  [10/17] Generative Agents")

    # 11. Day Pipeline
    slide_day_pipeline(prs)
    print("  [11/17] Day Pipeline")

    # 12. Night Pipeline
    slide_night_pipeline(prs)
    print("  [12/17] Night Pipeline")

    # 13. Dynamic Stats
    slide_dynamic_stats(prs)
    print("  [13/17] Dynamic Stats")

    # 14. Rule Engine
    slide_rule_engine(prs)
    print("  [14/17] Rule Engine")

    # 15. Section: Scenario
    slide_section_divider(prs, 5, "시나리오 상세", "코렐라인 — 인형이 된다는 것의 공포")
    print("  [15/17] Section divider")

    # 16. Characters
    slide_characters(prs)
    print("  [16/17] Characters")

    # 17. World Vars
    slide_world_vars(prs)
    print("  [17/20] World Vars")

    # 18. Endings
    slide_endings(prs)
    print("  [18/20] Endings")

    # 19. Dev Status
    slide_dev_status(prs)
    print("  [19/20] Dev Status")

    # 20. Future
    slide_future(prs)
    print("  [20/21] Future")

    # 21. Closing
    slide_closing(prs)
    print("  [21/21] Closing")

    output_path = r"C:\Users\jinhy\Jinhyeok Jeon\04 Projects\11 멋사_DeusExMachina\demo-repository\DeusExMachina_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✅ Saved: {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
